from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 增强配色方案
    primary_color = RGBColor(26, 35, 126)      # 深蓝
    accent_color = RGBColor(220, 53, 69)       # 红色强调
    secondary_color = RGBColor(52, 152, 219)   # 亮蓝
    text_color = RGBColor(33, 37, 41)          # 深灰
    light_bg = RGBColor(248, 249, 250)         # 浅灰
    gold_color = RGBColor(255, 193, 7)         # 金色
    
    # Slide 1: 封面 - 增强版
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 渐变效果背景（用两层模拟）
    bg1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = primary_color
    bg1.line.fill.background()
    
    # 装饰圆形
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(-1), Inches(4), Inches(4))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = secondary_color
    circle1.fill.transparency = 0.7
    circle1.line.fill.background()
    
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(5), Inches(3.5), Inches(3.5))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = accent_color
    circle2.fill.transparency = 0.7
    circle2.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "国文汇通商业计划书"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # 装饰线条 - 金色
    line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(4), Inches(5), Inches(0.08))
    line1.fill.solid()
    line1.fill.fore_color.rgb = gold_color
    line1.line.fill.background()
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "数字文创资产合规发行与交易一站式平台"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(26)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # 底部装饰
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "2026 · 汇聚文化之光 · 通达数字未来"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(16)
    footer_para.font.color.rgb = gold_color
    footer_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2-16: 内容页
    slides_content = [
        ("1. 政策东风：文化数字化的国家战略", [
            ("顶层设计：战略定调", [
                "2022年5月：中办、国办印发《关于推进实施国家文化数字化战略的意见》",
                "战略高度：为文化数字化发展提供了根本遵循，是国家级战略方向"
            ]),
            ("战略目标：两步走", [
                "2025年：基本建成文化数字化基础设施和服务平台",
                "2035年：建成国家文化大数据体系，中华文化全景呈现"
            ]),
            ("重点任务与机遇", [
                "任务：关联形成中华文化数据库、夯实基础设施",
                "机遇：明确'培育一批新型文化企业'，平台型企业迎来爆发期"
            ])
        ]),
        ("2. 行业浪潮：数字文创资产的崛起", [
            ("市场规模：万亿级蓝海", [
                "13.5万亿元：2024年中国数字文化行业市场规模",
                "3000亿元：2027年数字资产细分市场规模预测"
            ]),
            ("核心人群：Z世代主导", [
                "67%占比：Z世代（95后、00后）是绝对主力",
                "消费特征：追求个性化、情感共鸣，愿为文化价值付费"
            ]),
            ("行业三大趋势", [
                "合规化：监管完善，正规军进场，合规是生存前提",
                "数实融合：赋能实体经济，连接线上线下消费",
                "生态化：从单一交易走向'IP孵化-发行-应用'完整生态"
            ])
        ]),
        ("3. 痛点与解决方案：国文汇通的使命", [
            ("市场痛点", [
                "交易乱象：缺乏权威监管，投机炒作严重",
                "权益受损：确权不清，创作者和消费者权益无保障"
            ]),
            ("我们的解决方案", [
                "政策导向：完全响应国家战略，做合规排头兵",
                "技术赋能：利用区块链实现资产确权、可追溯",
                "重塑秩序：构建权威、透明的交易环境，让价值高效流转"
            ])
        ]),
        ("4. 核心护城河（一）：权威合规与渠道", [
            ("国资背书，权威合规", [
                "文交所直属：拥有国企背景，双重监管加持",
                "业务可溯：全流程透明，公信力远超行业竞品"
            ]),
            ("资质完备，技术领先", [
                "全牌照：持有全链条合规证照",
                "国家级链：接入星火链网子链（紫金山链），确权存证能力顶尖"
            ]),
            ("全渠道覆盖", [
                "APP全端上架：覆盖主流应用商店，获客通道畅通",
                "支付便捷：搭建顶尖支付通道，即将对接省级结算中心存管"
            ])
        ]),
        ("5. 核心护城河（二）：模型与团队", [
            ("经济模型：稳健增长", [
                "拒绝恶炒：以'寄存复利+生态多元'为核心",
                "价值锚定：绑定真实需求，资产价值稳步提升"
            ]),
            ("生态闭环：数实融合", [
                "闭环链路：数字获取 → 通宝流转 → 实体消费",
                "脱虚向实：依托五大模块，摆脱单纯的市场情绪驱动"
            ]),
            ("顶尖操盘团队", [
                "实战经验：团队曾操盘 Ibox、唯一艺术等顶级平台",
                "跨界资源：依托省文投，实现'体育+文旅'跨界布局"
            ])
        ]),
        ("6. 核心团队介绍", [
            ("战略负责人：佟世天", [
                "背景：带领'唯一艺术'原班人马",
                "职责：负责国文汇通战略规划，国内最早布局数字文化资产的专家"
            ]),
            ("运营负责人：'喜鹊'团队", [
                "背景：原 IBOX 平台总运营领衔",
                "战绩：曾推动IBOX市值超百亿、单日交易额近20亿",
                "职责：负责增长策略、用户体系与生态建设"
            ]),
            ("首席宣传官 (CMO)：Jeff", [
                "背景：链大教育创始人，Cardano 亚洲增长负责人",
                "职责：负责品牌战略、产业传播与全球化布局"
            ])
        ]),
        ("7. 平台战略资源矩阵", [
            ("交易所资源", [
                "深度绑定：江苏、四川、甘肃、安徽等省级文交所"
            ]),
            ("股东与IP资源", [
                "强力支撑：江苏文投、凤凰出版、省演艺集团",
                "内容供给：拥有海量优质文化IP授权"
            ]),
            ("地方文旅合作", [
                "标杆案例：重点合作福建泉州市文旅局",
                "模式复制：将地方文旅资源转化为国家级平台优势"
            ]),
            ("品牌合作", [
                "主要伙伴：泉州木偶剧院、紫金文创等，构建高质量内容库"
            ])
        ]),
        ("8. 生态模型：'一体两翼'", [
            ("一体：国家文化数据中枢", [
                "定位：连接国家文化大数据体系与市场",
                "功能：进行标准化的确权、登记、评估与发行"
            ]),
            ("左翼：标准制定", [
                "掌握话语权：制定文化数字资产全流程行业标准"
            ]),
            ("右翼：信任基座", [
                "技术互信：融合前沿技术，构建跨平台、可互信的技术底层"
            ])
        ]),
        ("9. 生态模型：'三维四域'", [
            ("三维体系", [
                "数据层：文博典藏·数字永生（文物高精度数字化，价值再生）",
                "平台层：非遗活化·产业赋能（提供交易市场与创意集市）",
                "场景层：文旅融合·虚实互促（嵌入商业空间，兑现价值）"
            ]),
            ("四域深耕", [
                "核心逻辑：从数据采集到场景应用，全链路深耕",
                "最终目标：打造'数字资产+线下消费'新模式"
            ])
        ]),
        ("10. 商业模式：通卷与通宝", [
            ("核心资产：国文通卷", [
                "总量：1000万份",
                "定位：平台基石资产，象征核心共建者权益",
                "获取：早期参与者持有"
            ]),
            ("流通积分：国文通宝", [
                "总量：10亿个",
                "通胀（产出）：用户寄存'通卷'，每日获得'通宝'奖励",
                "通缩（消耗）：用于交易手续费、打新、兑换商品",
                "价值逻辑：通宝在生态内流转，连接一切价值"
            ])
        ]),
        ("11. 经济模型：六大通缩机制", [
            ("1. 回收注销 (核心)", [
                "平台从收益中回购并永久注销通宝，减少流通量，提升价值"
            ]),
            ("2. 刚性消耗", [
                "手续费抵扣：交易必用",
                "资产发行：项目方需支付通宝作为认证费"
            ]),
            ("3. 场景消耗", [
                "企业入驻：品牌方质押通宝作为保证金",
                "实体兑换：直接兑换商品或服务",
                "生态博弈：趣味玩法消耗"
            ])
        ]),
        ("12. 发展路线图：三步走战略", [
            ("第一阶段：奠定基石 (1-2年)", [
                "策略：高通缩、低通胀",
                "目标：快速积累早期核心用户，完成平台冷启动，建立标杆"
            ]),
            ("第二阶段：生态稳健 (3-4年)", [
                "策略：平衡通胀通缩，精细化运营",
                "目标：验证商业模式，进入稳定发展期，扩大市场份额"
            ]),
            ("第三阶段：价值爆发 (5年以上)", [
                "策略：拓展应用场景，网络效应最大化",
                "目标：成为行业领导者，共享文化价值，实现全球化"
            ])
        ]),
        ("13. 用户多元化收益", [
            ("静态收益", [
                "寄存挖矿：持有资产，每日自动获得奖励"
            ]),
            ("动态收益", [
                "二级交易：藏品低买高卖，获取价差",
                "打新收益：抢购首发稀缺资产"
            ]),
            ("消费权益", [
                "实体赋能：凭藏品/通宝在合作商户享受折扣",
                "福利赠送：不定期获得文旅门票、演出票"
            ])
        ]),
        ("14. 五大应用场景落地", [
            ("1. 文博资源确权", [
                "博物馆馆藏文物上链，构建国家级文化资产库"
            ]),
            ("2. 非遗数字化", [
                "传统工艺、戏曲记录与资产化，支持传承人融资"
            ]),
            ("3. 数字文旅消费", [
                "景区门票、酒店权益NFT化，打通线上线下闭环"
            ]),
            ("4. 产业协同", [
                "开放API，支持第三方文旅集团接入；与政府共建城市级专区"
            ]),
            ("5. 文化出海", [
                "对接海外市场，推动中华文化数字资产全球流通"
            ])
        ]),
        ("15. 愿景与使命", [])
    ]
    
    # 生成内容页
    for idx, (title, content) in enumerate(slides_content, 2):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 白色背景
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg.line.fill.background()
        
        # 顶部装饰条
        header_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = primary_color
        header_bg.line.fill.background()
        
        # 左侧装饰条
        side_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.2), Inches(0.15), Inches(6.3))
        side_bar.fill.solid()
        side_bar.fill.fore_color.rgb = accent_color
        side_bar.line.fill.background()
        
        # 装饰圆点
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.3), Inches(0.4), Inches(0.4), Inches(0.4))
        dot.fill.solid()
        dot.fill.fore_color.rgb = gold_color
        dot.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.3), Inches(8.5), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(30)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        
        # 底部装饰线
        bottom_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.15), prs.slide_width, Inches(0.05))
        bottom_line.fill.solid()
        bottom_line.fill.fore_color.rgb = gold_color
        bottom_line.line.fill.background()
        
        # 页码
        page_box = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.3))
        page_frame = page_box.text_frame
        page_frame.text = f"{idx}/16"
        page_para = page_frame.paragraphs[0]
        page_para.font.size = Pt(12)
        page_para.font.color.rgb = text_color
        page_para.alignment = PP_ALIGN.RIGHT
        
        # 特殊处理最后一页（愿景与使命）
        if idx == 16:
            add_vision_mission_slide(slide, primary_color, accent_color, text_color, light_bg, gold_color)
        else:
            add_enhanced_content(slide, content, text_color, accent_color, secondary_color)
    
    prs.save('g:/PPT/国文汇通商业计划书_美化版.pptx')
    print("✨ 美化版PPT创建成功！")
    print("📁 文件保存在: g:/PPT/国文汇通商业计划书_美化版.pptx")

def add_enhanced_content(slide, content, text_color, accent_color, secondary_color):
    """添加增强的内容区域"""
    y_position = 1.5
    
    for section_title, items in content:
        # 章节标题卡片
        card_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                        Inches(0.4), Inches(y_position), 
                                        Inches(9.2), Inches(0.45))
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = RGBColor(240, 248, 255)
        card_bg.line.color.rgb = secondary_color
        card_bg.line.width = Pt(1.5)
        
        # 章节标题图标
        icon = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                     Inches(0.55), Inches(y_position + 0.08), 
                                     Inches(0.3), Inches(0.3))
        icon.fill.solid()
        icon.fill.fore_color.rgb = accent_color
        icon.line.fill.background()
        
        # 章节标题文字
        title_box = slide.shapes.add_textbox(Inches(1), Inches(y_position + 0.05), 
                                            Inches(8), Inches(0.35))
        title_frame = title_box.text_frame
        title_frame.text = section_title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(20)
        title_para.font.bold = True
        title_para.font.color.rgb = accent_color
        
        y_position += 0.55
        
        # 内容项
        for item in items:
            # 项目符号
            bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, 
                                          Inches(0.7), Inches(y_position + 0.05), 
                                          Inches(0.12), Inches(0.12))
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = secondary_color
            bullet.line.fill.background()
            
            # 内容文字
            item_box = slide.shapes.add_textbox(Inches(1), Inches(y_position), 
                                               Inches(8.5), Inches(0.3))
            item_frame = item_box.text_frame
            item_frame.text = item
            item_frame.word_wrap = True
            item_para = item_frame.paragraphs[0]
            item_para.font.size = Pt(14)
            item_para.font.color.rgb = text_color
            item_para.line_spacing = 1.3
            
            y_position += 0.35
        
        y_position += 0.15

def add_vision_mission_slide(slide, primary_color, accent_color, text_color, light_bg, gold_color):
    """添加愿景与使命页面"""
    y_pos = 1.8
    
    # 使命卡片
    mission_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                         Inches(1), Inches(y_pos), 
                                         Inches(8), Inches(1.6))
    mission_card.fill.solid()
    mission_card.fill.fore_color.rgb = light_bg
    mission_card.line.color.rgb = accent_color
    mission_card.line.width = Pt(3)
    
    # 使命图标
    mission_icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, 
                                         Inches(1.3), Inches(y_pos + 0.15), 
                                         Inches(0.4), Inches(0.4))
    mission_icon.fill.solid()
    mission_icon.fill.fore_color.rgb = accent_color
    mission_icon.line.fill.background()
    
    mission_frame = mission_card.text_frame
    mission_frame.margin_top = Inches(0.25)
    mission_frame.margin_left = Inches(0.8)
    
    p = mission_frame.paragraphs[0]
    p.text = "使命 MISSION"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = primary_color
    
    p = mission_frame.add_paragraph()
    p.text = "文化：让中华优秀传统文化在数字时代'活起来'"
    p.font.size = Pt(15)
    p.font.color.rgb = text_color
    p.space_before = Pt(8)
    
    p = mission_frame.add_paragraph()
    p.text = "产业：建立合规、安全、可持续的数字文化资产体系"
    p.font.size = Pt(15)
    p.font.color.rgb = text_color
    p.space_before = Pt(4)
    
    # 愿景卡片
    y_pos += 1.9
    vision_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                        Inches(1), Inches(y_pos), 
                                        Inches(8), Inches(1.6))
    vision_card.fill.solid()
    vision_card.fill.fore_color.rgb = light_bg
    vision_card.line.color.rgb = primary_color
    vision_card.line.width = Pt(3)
    
    # 愿景图标
    vision_icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, 
                                        Inches(1.3), Inches(y_pos + 0.15), 
                                        Inches(0.4), Inches(0.4))
    vision_icon.fill.solid()
    vision_icon.fill.fore_color.rgb = primary_color
    vision_icon.line.fill.background()
    
    vision_frame = vision_card.text_frame
    vision_frame.margin_top = Inches(0.25)
    vision_frame.margin_left = Inches(0.8)
    
    p = vision_frame.paragraphs[0]
    p.text = "愿景 VISION"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = primary_color
    
    p = vision_frame.add_paragraph()
    p.text = "近期：成为中国领先的文化类数字资产基础设施平台"
    p.font.size = Pt(15)
    p.font.color.rgb = text_color
    p.space_before = Pt(8)
    
    p = vision_frame.add_paragraph()
    p.text = "远期：构建全球领先的文化数字资产生态，通达数字未来"
    p.font.size = Pt(15)
    p.font.color.rgb = text_color
    p.space_before = Pt(4)
    
    # 底部口号
    y_pos += 2
    slogan_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(0.8))
    slogan_frame = slogan_box.text_frame
    slogan_frame.text = "汇聚文化之光 · 通达数字未来"
    slogan_para = slogan_frame.paragraphs[0]
    slogan_para.font.size = Pt(36)
    slogan_para.font.bold = True
    slogan_para.font.color.rgb = gold_color
    slogan_para.alignment = PP_ALIGN.CENTER

if __name__ == "__main__":
    create_presentation()
