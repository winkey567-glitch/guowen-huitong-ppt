from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 配色方案 (对应设计规则)
    primary_color = RGBColor(10, 10, 10)       # 深空黑 (背景)
    accent_color = RGBColor(255, 215, 0)       # 金色 (强调)
    secondary_color = RGBColor(26, 35, 126)    # 深蓝 (辅助)
    text_color = RGBColor(249, 249, 247)       # 浅灰 (文字)
    light_bg = RGBColor(249, 249, 247)         # 浅色背景(备用)
    
    # 字体设置 (注: PPT生成时字体依赖系统安装,这里设为通用字体)
    title_font = "Arial"
    body_font = "Arial"

    # Slide 1: 封面 (Rule: 全屏视觉冲击)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, primary_color)
    
    # 装饰圆
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2), Inches(1.5), Inches(6), Inches(6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = secondary_color
    circle.fill.transparency = 0.8
    circle.line.fill.background()

    add_title(slide, "国文汇通商业计划书", 60, accent_color, Inches(2.5))
    add_subtitle(slide, "数字文创资产合规发行与交易一站式平台", 24, text_color, Inches(4.5))
    
    # Slide 2: 政策东风 (Rule: 时间轴布局)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, primary_color)
    add_page_header(slide, "1. 政策东风：文化数字化的国家战略", accent_color)
    
    create_timeline(slide, [
        ("2022.05", "中办国办印发《关于推进实施国家文化数字化战略的意见》", "顶层设计"),
        ("2025", "基本建成文化数字化基础设施和服务平台", "基础设施"),
        ("2035", "建成国家文化大数据体系，中华文化全景呈现", "全面建成")
    ], text_color, accent_color)

    # Slide 3: 行业浪潮 (Rule: 数据可视化)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, primary_color)
    add_page_header(slide, "2. 行业浪潮：数字文创资产的崛起", accent_color)
    
    create_data_cards(slide, [
        ("13.5万亿", "2024年数字文化行业规模"),
        ("3000亿", "2027年数字资产市场预测"),
        ("67%", "Z世代(95后00后)消费占比")
    ], text_color, accent_color)
    
    # Slide 4: 国文汇通为何出现 (Rule: 三列布局)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, primary_color)
    add_page_header(slide, "3. 国文汇通为何出现？", accent_color)
    
    create_3col_layout(slide, [
        ("市场痛点", "交易乱象、权益受损", "亟需权威合规平台"),
        ("政策导向", "国家战略、数字中国", "做合规排头兵"),
        ("技术赋能", "区块链确权、可追溯", "让价值高效流转")
    ], text_color, accent_color)

    # Slide 5: 核心护城河一 (Rule: 四象限/网格)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, primary_color)
    add_page_header(slide, "4. 核心护城河(一)：资质与渠道", accent_color)
    
    create_grid_layout(slide, [
        ("权威合规", "文交所直属，国企背景，双重监管"),
        ("资质完备", "全牌照，接入星火链网(紫金山链)"),
        ("全渠道", "APP全端上架，主流应用商店覆盖"),
        ("支付便捷", "顶尖支付通道，省级结算中心存管")
    ], text_color, accent_color)

    # Slide 6: 核心护城河二 (Rule: 流程图逻辑)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, primary_color)
    add_page_header(slide, "5. 核心护城河(二)：模型与团队", accent_color)
    
    create_list_content(slide, [
        ("经济模型", "寄存复利 + 生态多元，拒绝恶炒"),
        ("生态闭环", "数字获取 -> 通宝流转 -> 实体消费"),
        ("顶尖团队", "操盘过Ibox、唯一艺术，跨界资源丰富")
    ], text_color)

    # Slide 7: 核心团队 (Rule: 人物卡片)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, primary_color)
    add_page_header(slide, "6. 核心团队介绍", accent_color)
    
    create_3col_layout(slide, [
        ("佟世天", "战略负责人", "带领'唯一艺术'原班人马，国内最早布局者"),
        ("喜鹊团队", "运营负责人", "原IBOX总运营，曾推动百亿市值"),
        ("Jeff", "CMO", "链大教育创始人，Cardano亚洲增长负责人")
    ], text_color, accent_color)

    # Slide 8: 战略资源 (Rule: 列表/地图)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, primary_color)
    add_page_header(slide, "7. 平台战略资源矩阵", accent_color)
    
    create_grid_layout(slide, [
        ("交易所", "江苏/四川/甘肃/安徽文交所深度绑定"),
        ("股东/IP", "江苏文投、凤凰出版、省演艺集团"),
        ("文旅合作", "福建泉州文旅局标杆案例"),
        ("品牌伙伴", "泉州木偶剧院、紫金文创等")
    ], text_color, accent_color)

    # Slide 9: 生态模型一 (Rule: 架构图)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, primary_color)
    add_page_header(slide, "8. 生态模型：一体两翼", accent_color)
    
    create_list_content(slide, [
        ("一体：国家文化数据中枢", "连接国家文化大数据体系与市场"),
        ("左翼：标准制定", "制定全流程行业标准，掌握话语权"),
        ("右翼：信任基座", "构建跨平台、可互信的技术底层")
    ], text_color)

    # Slide 10: 生态模型二 (Rule: 分层)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, primary_color)
    add_page_header(slide, "9. 生态模型：三维四域", accent_color)
    
    create_list_content(slide, [
        ("三维体系", "数据层(珍资源) - 平台层(交易) - 场景层(消费)"),
        ("四域深耕", "文博典藏、非遗活化、文旅融合、产业赋能")
    ], text_color)

    # Slide 11: 商业模式 (Rule: 双币模型)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, primary_color)
    add_page_header(slide, "10. 商业模式：通卷与通宝", accent_color)
    
    create_2col_layout(slide, 
        ("国文通卷 (1000万份)", "平台基石资产\n权益象征\n早期持有"),
        ("国文通宝 (10亿个)", "流通积分\n寄存产出\n消费/打新消耗"),
        text_color, accent_color
    )

    # Slide 12: 经济模型 (Rule: 六大通缩)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, primary_color)
    add_page_header(slide, "11. 经济模型：六大通缩机制", accent_color)
    
    create_grid_layout(slide, [
        ("回收注销", "平台回购并永久注销"),
        ("手续费抵扣", "交易必用，刚性消耗"),
        ("资产发行", "项目方支付通宝作为认证费"),
        ("企业入驻", "品牌方质押保证金"),
        ("实体兑换", "直接兑换商品服务"),
        ("生态博弈", "趣味玩法消耗")
    ], text_color, accent_color, columns=3)

    # Slide 13: 动态平衡 (Rule: 平衡图)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, primary_color)
    add_page_header(slide, "12. 经济模型：动态平衡", accent_color)
    
    create_2col_layout(slide,
        ("通胀 (注入)", "通卷寄存挖矿\n每日通宝产出\n为生态提供动力"),
        ("通缩 (捕获)", "回收注销\n手续费/发行认证\n生态消费消耗"),
        text_color, accent_color
    )

    # Slide 14: 发展路线 (Rule: 路线图)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, primary_color)
    add_page_header(slide, "13. 发展路线图：三步走", accent_color)
    
    create_timeline(slide, [
        ("第一阶段", "奠定基石 (1-2年)", "高通缩低通胀，积累核心用户"),
        ("第二阶段", "生态稳健 (3-4年)", "平衡通胀通缩，精细化运营"),
        ("第三阶段", "价值爆发 (5年以上)", "拓展应用场景，全球化")
    ], text_color, accent_color)

    # Slide 15: 用户收益 (Rule: Bento Grid)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, primary_color)
    add_page_header(slide, "14. 用户多元化收益", accent_color)
    
    create_grid_layout(slide, [
        ("静态收益", "寄存挖矿获通宝"),
        ("动态收益", "二级交易价差"),
        ("打新收益", "抢购首发稀缺资产"),
        ("消费权益", "实体折扣与福利")
    ], text_color, accent_color)

    # Slide 16-20: 应用场景 (Rule: 场景展示)
    scenes = [
        ("15. 场景一：文化资产确权", "文博资源上链，构建国家级库；非遗数字化融资"),
        ("16. 场景二：数字文旅消费", "景区门票/酒店权益NFT化，打通线上线下"),
        ("17. 场景三：互动娱乐", "GameFi(边玩边赚) + StudyFi(学习激励)"),
        ("18. 场景四：产业协同", "开放API给第三方，与政府共建城市专区"),
        ("19. 场景五：文化出海", "对接海外平台，多语言版本，全球流通")
    ]
    
    for title, content in scenes:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide, primary_color)
        add_page_header(slide, title, accent_color)
        
        box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3))
        tf = box.text_frame
        tf.text = content
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.color.rgb = text_color
        tf.word_wrap = True

    # Slide 21: 战略规划
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, primary_color)
    add_page_header(slide, "20. 战略规划：未来蓝图", accent_color)
    create_timeline(slide, [
        ("短期", "1-2年", "标杆打造，模式验证"),
        ("中期", "3-4年", "生态扩张，网络效应"),
        ("长期", "5年以上", "全球影响，标准输出")
    ], text_color, accent_color)

    # Slide 22: 使命
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, primary_color)
    add_center_text(slide, "使命 MISSION", "让中华优秀传统文化在数字时代'活起来'\n建立合规、安全、可持续的数字文化资产体系", accent_color, text_color)

    # Slide 23: 愿景
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, primary_color)
    add_center_text(slide, "愿景 VISION", "汇聚文化之光 · 通达数字未来\n成为中国领先的文化类数字资产基础设施平台", accent_color, text_color)

    # Slide 24: 封底
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, primary_color)
    add_center_text(slide, "感谢聆听", "国文汇通\nBusiness Plan 2026", accent_color, text_color)

    prs.save('g:/PPT/国文汇通商业计划书_24页完整版.pptx')
    print("✨ 24页完整版PPT创建成功！")

# Helper Functions

def set_background(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()

def add_title(slide, text, size, color, top):
    box = slide.shapes.add_textbox(Inches(1), top, Inches(8), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

def add_subtitle(slide, text, size, color, top):
    add_title(slide, text, size, color, top)

def add_page_header(slide, text, color):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = color
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(9), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = color

def create_timeline(slide, items, text_color, accent_color):
    left = 1
    width = 8 / len(items)
    for time, title, desc in items:
        # Time
        box = slide.shapes.add_textbox(Inches(left), Inches(2.5), Inches(width-0.2), Inches(0.5))
        p = box.text_frame.paragraphs[0]
        p.text = time
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = accent_color
        
        # Dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(3.2), Inches(0.2), Inches(0.2))
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent_color
        
        # Title
        box = slide.shapes.add_textbox(Inches(left), Inches(3.5), Inches(width-0.2), Inches(0.8))
        p = box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = text_color
        
        # Desc
        box = slide.shapes.add_textbox(Inches(left), Inches(4.5), Inches(width-0.2), Inches(2))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = text_color
        
        left += width

def create_data_cards(slide, items, text_color, accent_color):
    left = 1
    width = 8 / len(items)
    for value, desc in items:
        box = slide.shapes.add_textbox(Inches(left), Inches(3), Inches(width-0.2), Inches(1))
        p = box.text_frame.paragraphs[0]
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = accent_color
        p.alignment = PP_ALIGN.CENTER
        
        box = slide.shapes.add_textbox(Inches(left), Inches(4.2), Inches(width-0.2), Inches(1))
        p = box.text_frame.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER
        
        left += width

def create_3col_layout(slide, items, text_color, accent_color):
    left = 0.5
    width = 3
    for title, sub, desc in items:
        # Card bg
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2), Inches(2.8), Inches(4.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(30, 30, 30)
        bg.line.color.rgb = accent_color
        bg.line.width = Pt(1)
        
        # Title
        box = slide.shapes.add_textbox(Inches(left+0.2), Inches(2.2), Inches(2.4), Inches(0.5))
        p = box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = accent_color
        
        # Sub
        box = slide.shapes.add_textbox(Inches(left+0.2), Inches(2.8), Inches(2.4), Inches(0.5))
        p = box.text_frame.paragraphs[0]
        p.text = sub
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = text_color
        
        # Desc
        box = slide.shapes.add_textbox(Inches(left+0.2), Inches(3.5), Inches(2.4), Inches(2.5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = text_color
        
        left += 3.1

def create_grid_layout(slide, items, text_color, accent_color, columns=2):
    row = 0
    col = 0
    start_x = 1
    start_y = 2
    w = 8 / columns
    h = 2
    
    for title, desc in items:
        x = start_x + col * w
        y = start_y + row * h
        
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w-0.2), Inches(0.4))
        p = box.text_frame.paragraphs[0]
        p.text = "• " + title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = accent_color
        
        box = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+0.5), Inches(w-0.4), Inches(1.2))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = text_color
        
        col += 1
        if col >= columns:
            col = 0
            row += 1

def create_list_content(slide, items, text_color):
    top = 2
    for title, desc in items:
        box = slide.shapes.add_textbox(Inches(1), Inches(top), Inches(8), Inches(0.4))
        p = box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 215, 0)
        
        box = slide.shapes.add_textbox(Inches(1.5), Inches(top+0.5), Inches(7.5), Inches(0.8))
        p = box.text_frame.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = text_color
        
        top += 1.5

def create_2col_layout(slide, left_data, right_data, text_color, accent_color):
    # Left
    l_title, l_desc = left_data
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3.5), Inches(0.5))
    p = box.text_frame.paragraphs[0]
    p.text = l_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = accent_color
    
    box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(3.5), Inches(3))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = l_desc
    p.font.size = Pt(14)
    p.font.color.rgb = text_color
    
    # Right
    r_title, r_desc = right_data
    box = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(3.5), Inches(0.5))
    p = box.text_frame.paragraphs[0]
    p.text = r_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = accent_color
    
    box = slide.shapes.add_textbox(Inches(5.5), Inches(2.8), Inches(3.5), Inches(3))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = r_desc
    p.font.size = Pt(14)
    p.font.color.rgb = text_color

def add_center_text(slide, title, subtitle, accent_color, text_color):
    box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = accent_color
    p.alignment = PP_ALIGN.CENTER
    
    box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
    p = box.text_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(18)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

if __name__ == "__main__":
    create_presentation()
