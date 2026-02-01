# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

pc = RGBColor(26, 35, 126)
ac = RGBColor(220, 53, 69)
tc = RGBColor(33, 37, 41)
lb = RGBColor(248, 249, 250)
wc = RGBColor(255, 255, 255)

def add_header(slide, title):
    hbg = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(1))
    hbg.fill.solid()
    hbg.fill.fore_color.rgb = pc
    hbg.line.fill.background()
    
    htb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    htf = htb.text_frame
    htf.text = title
    hp = htf.paragraphs[0]
    hp.font.size = Pt(28)
    hp.font.bold = True
    hp.font.color.rgb = wc
    
    hln = slide.shapes.add_shape(1, 0, Inches(1), Inches(10), Inches(0.05))
    hln.fill.solid()
    hln.fill.fore_color.rgb = ac
    hln.line.fill.background()

def add_bullet(slide, y, text):
    itb = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(8), Inches(0.3))
    itf = itb.text_frame
    itf.text = text
    itp = itf.paragraphs[0]
    itp.font.size = Pt(14)
    itp.font.color.rgb = tc
    return y + 0.3

slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = pc
bg.line.fill.background()

tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
tf = tb.text_frame
tf.text = u'国文汇通商业计划书'
p = tf.paragraphs[0]
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = wc
p.alignment = PP_ALIGN.CENTER

stb = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.8))
stf = stb.text_frame
stf.text = u'数字文创资产合规发行与交易一站式平台'
sp = stf.paragraphs[0]
sp.font.size = Pt(24)
sp.font.color.rgb = wc
sp.alignment = PP_ALIGN.CENTER

ln = slide.shapes.add_shape(1, Inches(3), Inches(5.3), Inches(4), Inches(0.05))
ln.fill.solid()
ln.fill.fore_color.rgb = ac
ln.line.fill.background()

for i in range(2, 17):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, f'Slide {i}')
    y = 1.5
    y = add_bullet(slide, y, f'Content for slide {i}')

prs.save('GuowenHuitong_Final.pptx')
print('PPT created with 16 slides!')
