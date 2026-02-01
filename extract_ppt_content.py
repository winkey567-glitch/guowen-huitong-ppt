"""
提取 PowerPoint 文件的文字内容
"""
from pptx import Presentation
import json

def extract_ppt_content(ppt_path):
    """提取PPT中的所有文字内容"""
    prs = Presentation(ppt_path)
    
    slides_content = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_data = {
            'slide_number': slide_num,
            'title': '',
            'content': []
        }
        
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    # 尝试识别标题
                    if shape.is_placeholder and shape.placeholder_format.type == 1:
                        slide_data['title'] = text
                    else:
                        slide_data['content'].append(text)
        
        slides_content.append(slide_data)
    
    return slides_content

if __name__ == "__main__":
    ppt_path = "国文汇通商业计划书（整理v1.1）.pptx"
    
    try:
        content = extract_ppt_content(ppt_path)
        
        # 保存为 JSON
        with open('ppt_content.json', 'w', encoding='utf-8') as f:
            json.dump(content, f, ensure_ascii=False, indent=2)
        
        # 保存为 Markdown
        with open('ppt_content.md', 'w', encoding='utf-8') as f:
            for slide in content:
                f.write(f"\n## Slide {slide['slide_number']}\n\n")
                if slide['title']:
                    f.write(f"### {slide['title']}\n\n")
                for text in slide['content']:
                    f.write(f"{text}\n\n")
        
        print(f"✅ 成功提取 {len(content)} 页内容")
        print(f"📄 已保存到: ppt_content.json 和 ppt_content.md")
        
    except Exception as e:
        print(f"❌ 错误: {e}")
        print("\n请确保已安装 python-pptx:")
        print("pip install python-pptx")
